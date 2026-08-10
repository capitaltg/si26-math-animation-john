import itertools
from dataclasses import dataclass
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

#: Ceiling on nested-group depth. Real decks nest a handful of levels; a
#: crafted archive with hundreds of GroupShape wrappers would blow the Python
#: recursion limit before python-pptx could return.
_MAX_GROUP_DEPTH = 32


@dataclass(frozen=True)
class Block:
    kind: Literal["text", "cell"]
    table_ord: int | None
    text: str


def _extract_shape_blocks(
    shape, table_ords: "itertools.count", depth: int = 0
) -> list[Block]:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if depth >= _MAX_GROUP_DEPTH:
            return []
        return [
            block
            for child in shape.shapes
            for block in _extract_shape_blocks(child, table_ords, depth + 1)
        ]
    if getattr(shape, "has_table", False):
        table_ord = next(table_ords)
        return [
            Block(kind="cell", table_ord=table_ord, text=cell.text)
            for row in shape.table.rows
            for cell in row.cells
        ]
    if getattr(shape, "has_text_frame", False):
        return [Block(kind="text", table_ord=None, text=shape.text_frame.text)]
    return []


def extract_slide_blocks(pptx_path: Path) -> list[list[Block]]:
    presentation = Presentation(pptx_path)
    slides_blocks: list[list[Block]] = []
    for slide in presentation.slides:
        table_ords = itertools.count()
        blocks: list[Block] = []
        for shape in slide.shapes:
            blocks.extend(_extract_shape_blocks(shape, table_ords))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            blocks.append(
                Block(kind="text", table_ord=None, text=slide.notes_slide.notes_text_frame.text)
            )
        slides_blocks.append([block for block in blocks if block.text.strip()])
    return slides_blocks


def flatten_blocks(blocks: list[Block]) -> str:
    return "\n".join(block.text for block in blocks if block.text.strip())


def chunk_slide_blocks(
    slide_blocks: list[list[Block]], chunk_size: int = 25
) -> list[list[list[Block]]]:
    return [slide_blocks[i:i + chunk_size] for i in range(0, len(slide_blocks), chunk_size)]
