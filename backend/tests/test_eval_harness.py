import sys
from pathlib import Path
from unittest.mock import patch

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))


def test_fixture_builders_write_all_four_decks(tmp_path):
    from pptx import Presentation

    from eval.generate_fixtures import (
        build_ambiguous_phrasing_deck,
        build_chain_test_deck,
        build_distractor_heavy_deck,
        build_zero_candidate_deck,
    )

    builders = [
        build_zero_candidate_deck,
        build_distractor_heavy_deck,
        build_ambiguous_phrasing_deck,
        build_chain_test_deck,
    ]
    paths = [tmp_path / f"fixture_{index}.pptx" for index in range(4)]

    for builder, path in zip(builders, paths):
        builder(path)

    assert all(path.exists() and path.stat().st_size > 0 for path in paths)

    chain_deck = Presentation(paths[-1])
    assert len(chain_deck.slides) == 3
    assert [
        (slide.shapes.title.text, slide.placeholders[1].text)
        for slide in chain_deck.slides
    ] == [
        (
            "Problem 1",
            "A frog is sitting on 3. It jumps forward 4 spaces. Where does it land?",
        ),
        (
            "Problem 2",
            "Start at 7. Move back 5. What number are you on?",
        ),
        (
            "Problem 3",
            "A snail starts at 2. It crawls forward 6. What number does it reach?",
        ),
    ]


@patch("eval.run_eval.discover_candidates_for_document", return_value=[])
def test_run_fixture_reports_zero_candidates(mock_discover, tmp_path):
    from eval.generate_fixtures import build_zero_candidate_deck
    from eval.run_eval import run_fixture

    pptx_path = tmp_path / "zero_candidate_deck.pptx"
    build_zero_candidate_deck(pptx_path)

    report = run_fixture(pptx_path)

    assert report == {
        "fixture": "zero_candidate_deck.pptx",
        "candidate_count": 0,
        "candidates": [],
    }
    mock_discover.assert_called_once()
