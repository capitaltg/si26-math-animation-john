from pptx import Presentation
from pptx.util import Inches


def _build_sample_pptx(path):
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "Warm Up"
    slide1.placeholders[1].text = "Sarah has 4 apples and buys 3 more. How many now?"
    slide1.notes_slide.notes_text_frame.text = "Remind students this is simple addition."

    table = slide1.shapes.add_table(
        1, 1, Inches(1), Inches(4), Inches(4), Inches(1)
    ).table
    table.cell(0, 0).text = "Table problem: 6 groups of 4"

    group = slide1.shapes.add_group_shape()
    nested_group = group.shapes.add_group_shape()
    grouped_text = nested_group.shapes.add_textbox(
        Inches(1), Inches(5), Inches(4), Inches(1)
    )
    grouped_text.text_frame.text = "Grouped problem: 9 minus 2"

    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Agenda"
    slide2.placeholders[1].text = "Standards: 3.OA.A.1"

    presentation.save(path)


def _build_two_table_pptx(path):
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Two tables"

    table_a = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    table_a.cell(0, 0).text = "6"
    table_a.cell(0, 1).text = "3"

    table_b = slide.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    table_b.cell(0, 0).text = "10"
    table_b.cell(0, 1).text = "20"

    presentation.save(path)


def test_extract_slide_blocks_tags_text_and_cell_blocks(tmp_path):
    from app.pipeline.parsing import extract_slide_blocks

    pptx_path = tmp_path / "sample.pptx"
    _build_sample_pptx(pptx_path)

    slide_blocks = extract_slide_blocks(pptx_path)

    assert len(slide_blocks) == 2
    slide1_blocks = slide_blocks[0]

    text_blocks = [b for b in slide1_blocks if b.kind == "text"]
    cell_blocks = [b for b in slide1_blocks if b.kind == "cell"]

    assert any("Sarah has 4 apples" in b.text for b in text_blocks)
    assert any("simple addition" in b.text for b in text_blocks)
    assert any("Grouped problem: 9 minus 2" in b.text for b in text_blocks)
    assert all(b.table_ord is None for b in text_blocks)

    assert len(cell_blocks) == 1
    assert cell_blocks[0].text == "Table problem: 6 groups of 4"
    assert cell_blocks[0].table_ord == 0

    slide2_blocks = slide_blocks[1]
    assert any("3.OA.A.1" in b.text for b in slide2_blocks)


def test_extract_slide_blocks_assigns_distinct_table_ord_per_table(tmp_path):
    from app.pipeline.parsing import extract_slide_blocks

    pptx_path = tmp_path / "two_tables.pptx"
    _build_two_table_pptx(pptx_path)

    slide_blocks = extract_slide_blocks(pptx_path)
    cell_blocks = [b for b in slide_blocks[0] if b.kind == "cell"]

    table_ords = {b.table_ord for b in cell_blocks}
    assert table_ords == {0, 1}
    assert {b.text for b in cell_blocks if b.table_ord == 0} == {"6", "3"}
    assert {b.text for b in cell_blocks if b.table_ord == 1} == {"10", "20"}


def test_extract_slide_blocks_assigns_table_ord_to_table_nested_in_group(tmp_path):
    from app.pipeline.parsing import extract_slide_blocks

    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Nested table"

    group = slide.shapes.add_group_shape()
    nested_group = group.shapes.add_group_shape()

    # python-pptx's GroupShapes has no add_table; build the table on the
    # slide's own shape tree, then relocate its XML element into the
    # nested group so it behaves as a table genuinely nested in a group.
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "7"
    table_shape.table.cell(0, 1).text = "8"
    table_shape._element.getparent().remove(table_shape._element)
    nested_group.shapes._spTree.append(table_shape._element)

    pptx_path = tmp_path / "nested_table.pptx"
    presentation.save(pptx_path)

    slide_blocks = extract_slide_blocks(pptx_path)
    cell_blocks = [b for b in slide_blocks[0] if b.kind == "cell"]

    assert len(cell_blocks) == 2
    assert {b.table_ord for b in cell_blocks} == {0}
    assert {b.text for b in cell_blocks} == {"7", "8"}


def test_flatten_blocks_matches_legacy_join_order(tmp_path):
    from app.pipeline.parsing import extract_slide_blocks, flatten_blocks

    pptx_path = tmp_path / "sample.pptx"
    _build_sample_pptx(pptx_path)

    slide_blocks = extract_slide_blocks(pptx_path)
    flat = flatten_blocks(slide_blocks[0])

    assert flat == "\n".join([
        "Warm Up",
        "Sarah has 4 apples and buys 3 more. How many now?",
        "Table problem: 6 groups of 4",
        "Grouped problem: 9 minus 2",
        "Remind students this is simple addition.",
    ])


def test_chunk_slide_blocks_splits_at_chunk_size():
    from app.pipeline.parsing import Block, chunk_slide_blocks

    slide_blocks = [[Block(kind="text", table_ord=None, text=f"slide {i}")] for i in range(50)]

    chunks = chunk_slide_blocks(slide_blocks, chunk_size=25)

    assert len(chunks) == 2
    assert len(chunks[0]) == 25
    assert len(chunks[1]) == 25
    assert chunks[0][0][0].text == "slide 0"
    assert chunks[1][0][0].text == "slide 25"
