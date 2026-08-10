import io
import zipfile
from pathlib import Path

import pytest

from app.pipeline.pptx_guard import (
    MAX_COMPRESSION_RATIO,
    MAX_MEMBER_COUNT,
    MAX_MEMBER_UNCOMPRESSED,
    MAX_UNCOMPRESSED_BYTES,
    PptxGuardError,
    inspect_pptx_archive,
)


def _write_zip(path: Path, members: list[tuple[str, bytes]], *, encrypt_first: bool = False) -> None:
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name, data in members:
            info = zipfile.ZipInfo(name)
            info.compress_type = zipfile.ZIP_DEFLATED
            if encrypt_first and name == members[0][0]:
                info.flag_bits |= 0x1
            zf.writestr(info, data)


def test_inspect_accepts_normal_pptx(tmp_path):
    from pptx import Presentation

    path = tmp_path / "ok.pptx"
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Hello"
    presentation.save(path)

    inspect_pptx_archive(path)


def test_inspect_rejects_non_zip(tmp_path):
    path = tmp_path / "junk.pptx"
    path.write_bytes(b"not a zip at all")

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason == "not_a_zip"


def test_inspect_rejects_zip_bomb_by_ratio(tmp_path):
    path = tmp_path / "bomb.pptx"
    # Highly compressible payload: 20 MB of zeros deflates to ~20 KB → ratio ~1000
    payload_size = 20 * 1024 * 1024
    _write_zip(path, [("bomb.xml", b"\x00" * payload_size)])

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    # Either ratio or per-member cap catches it first — both are legitimate.
    assert exc.value.reason in {"ratio_too_high", "member_too_large", "total_too_large"}


def test_inspect_rejects_huge_uncompressed_total(tmp_path):
    path = tmp_path / "huge.pptx"
    # Force TOTAL cap without tripping per-member or ratio: many members
    # of moderate size, low compressibility.
    import os
    per_member = MAX_MEMBER_UNCOMPRESSED - 1
    # Random-ish content resists deflate → ratio stays low.
    payload = os.urandom(1024)  # 1 KB seed
    filler = payload * (per_member // len(payload))
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as zf:
        for i in range(3):
            zf.writestr(f"part{i}.bin", filler)

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason in {"total_too_large", "member_too_large"}


def test_inspect_rejects_excessive_members(tmp_path):
    path = tmp_path / "many.pptx"
    with zipfile.ZipFile(path, "w") as zf:
        for i in range(MAX_MEMBER_COUNT + 5):
            zf.writestr(f"m{i}.xml", b"")

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason == "too_many_members"


def test_inspect_rejects_encrypted_member(tmp_path):
    path = tmp_path / "enc.pptx"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("secret.xml", b"payload")
    # zipfile.writestr does not honor a pre-set encryption flag; flip the
    # bit directly in the local header (offset 6) and central directory
    # entry (offset 8) so the archive looks encrypted to any zip reader.
    data = bytearray(path.read_bytes())
    lh = data.index(b"PK\x03\x04")
    data[lh + 6] |= 0x01
    cd = data.index(b"PK\x01\x02")
    data[cd + 8] |= 0x01
    path.write_bytes(bytes(data))

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason == "encrypted"


@pytest.mark.parametrize(
    "bad_name",
    [
        "../evil.xml",
        "a/../../etc/passwd",
        "/absolute.xml",
        "C:/windows/system32",
    ],
)
def test_inspect_rejects_suspicious_paths(tmp_path, bad_name):
    path = tmp_path / "trav.pptx"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr(bad_name, b"x")

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason == "suspicious_path"


def test_inspect_rejects_single_member_over_cap(tmp_path):
    path = tmp_path / "big_member.pptx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as zf:
        # Fake declared file_size by writing exactly the threshold + 1 bytes
        # of low-compressibility content (STORED so file_size==compress_size).
        import os
        chunk = os.urandom(1024)
        needed = MAX_MEMBER_UNCOMPRESSED + 1
        with zf.open(zipfile.ZipInfo("giant.bin"), "w") as fh:
            written = 0
            while written < needed:
                take = min(len(chunk), needed - written)
                fh.write(chunk[:take])
                written += take

    with pytest.raises(PptxGuardError) as exc:
        inspect_pptx_archive(path)
    assert exc.value.reason == "member_too_large"
