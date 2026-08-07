"""The teacher-facing half of the meta-template loop.

These routes are authorized by the ordinary session cookie rather than the
reviewer bearer token: the whole point is that a teacher can watch and approve
the template built from their own problem without an admin. Ownership is
therefore the only thing standing between one session's draft and another's, so
most of this file is about that boundary.
"""

import io
import json
from datetime import datetime, timedelta, timezone
from unittest.mock import patch
from uuid import uuid4

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app.config import get_settings
from app.meta import db, models
from app.meta.artifacts import store_artifact
from app.meta.draft_generation import DraftProposal, ProposedFixture
from app.meta.drafts import persist_reviewable_draft
from app.meta.dsl.expression import FieldRefNode
from app.meta.dsl.guard import GuardDocument, PositivePredicate
from app.meta.dsl.params import IntegerFieldSpec, ParamsDocument
from app.meta.dsl.teaching_plan import TeachingPlanDocument
from app.meta.dsl.v3_common import CompileContext
from app.meta.fingerprint import Fingerprint, canonical_fingerprint_key
from app.meta.v3.compiler import compile_teaching_plan
from app.meta.v3.quality import validate_static_quality
from app.meta.validation import FixtureCheckResult
from app.meta.validation_pipeline import ValidatedCandidate
from app.meta.versions import DSL_COMPILER_VERSION, DYNAMIC_RENDERER_VERSION


def _now():
    return datetime(2026, 8, 4, tzinfo=timezone.utc)


def _enable(monkeypatch, tmp_path, *, engine, dynamic_classifier=True, approval=True):
    monkeypatch.setattr(db, "get_engine", lambda: engine)
    monkeypatch.setenv("META_ARTIFACT_ROOT", str(tmp_path / "artifacts"))
    monkeypatch.setenv("META_TEMPLATES_ENABLED", "1")
    monkeypatch.setenv("META_CODEGEN_ENABLED", "1")
    monkeypatch.setenv("META_APPROVAL_ENABLED", "1" if approval else "0")
    monkeypatch.setenv("META_DYNAMIC_CLASSIFIER_ENABLED", "1" if dynamic_classifier else "0")
    monkeypatch.setenv("META_REVIEWER_TOKEN", "test-token")
    get_settings.cache_clear()


@pytest.fixture
def engine(tmp_path):
    eng = db.make_engine(tmp_path / "meta.db")
    db.create_all(eng)
    return eng


@pytest.fixture
def client(tmp_path, monkeypatch, engine):
    _enable(monkeypatch, tmp_path, engine=engine)
    from app.main import create_app

    yield TestClient(create_app())
    get_settings.cache_clear()


def _pptx_bytes():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide"
    slide.placeholders[1].text = "Which pair of socks is left over?"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _start_session(client, candidate_id="c1"):
    """Give the client a real session cookie with one candidate in it."""
    from app.models.candidate import Candidate

    candidate = Candidate(
        candidate_id=candidate_id,
        source_excerpt="Seven socks are laid out; which one has no pair?",
        slide_index=0,
        one_line_summary="Detected: odd one out",
    )
    with patch("app.routes.discover_candidates_for_document", return_value=[candidate]):
        client.post(
            "/upload",
            files={
                "file": (
                    "deck.pptx",
                    _pptx_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
    return client.cookies.get("session_id")


def _fingerprint():
    return Fingerprint(
        fingerprint_version=1,
        operation_family="compare",
        representation_family="set",
        number_domain="whole",
        operand_arity=7,
        step_count=1,
        grade_band="3-5",
    )


def _plan():
    return TeachingPlanDocument.model_validate({
        "plan_version": 3,
        "learning_objective": "Find the value with no partner.",
        "primary_visual": {"kind": "label", "ref": "n_label", "text": "value"},
        "strategy": "group_reveal",
        "beats": [
            {"id": "reveal", "kind": "reveal", "targets": [{"visual_ref": "n_label"}],
             "intent": "show the values together"},
            {"id": "focus", "kind": "focus", "targets": [{"visual_ref": "n_label"}],
             "intent": "focus on the unpaired value"},
            {"id": "conclude", "kind": "conclude", "targets": [{"visual_ref": "n_label"}],
             "intent": "state the leftover"},
        ],
        "variation_seed": "teacher-api",
    })


def _proposal(observation_id):
    return DraftProposal(
        params_document=ParamsDocument(
            params_version=1,
            fields=[IntegerFieldSpec(name="n", label="N", description="", minimum=1, maximum=10)],
        ),
        guard_document=GuardDocument(
            guard_version=1, predicates=[PositivePredicate(value=FieldRefNode(field="n"))]
        ),
        answer_expression=FieldRefNode(field="n"),
        teaching_plan_document=_plan(),
        classifier_bullet="use for leftover-pair problems",
        fixtures=[
            ProposedFixture(
                kind="positive", expected_outcome="accept",
                observation_id=observation_id, params={"n": 5},
            ),
            ProposedFixture(kind="negative", expected_outcome="reject", params={"n": -1}),
        ],
    )


def _candidate_bundle(proposal):
    scene_program = compile_teaching_plan(
        proposal.teaching_plan_document,
        proposal.answer_expression,
        frozenset({"n"}),
        CompileContext(concept_family="compare_collection", grade_band="3-5"),
    )
    quality = validate_static_quality(
        proposal.teaching_plan_document, scene_program
    ).model_payload()
    artifact_hash = f"sha256:{uuid4().hex}"
    quality["artifact_hash"] = artifact_hash
    preview_hash = store_artifact(
        get_settings().meta_artifact_root, b"fake preview bytes for teacher-api tests"
    )
    return ValidatedCandidate(
        proposal=proposal,
        scene_program=scene_program,
        validation_report={
            "passed": True,
            "fixture_results": [],
            "preview_ok": True,
            "preview_artifact_hash": preview_hash,
            "artifact_hash": artifact_hash,
            "compiler_version": DSL_COMPILER_VERSION,
            "renderer_version": DYNAMIC_RENDERER_VERSION,
            "negative_predicate_coverage": [0],
        },
        quality_report=quality,
        preview_artifact_hash=preview_hash,
        fixture_results=[
            FixtureCheckResult("fixture-0", True, "accepted"),
            FixtureCheckResult("fixture-1", True, "rejected", frozenset({0})),
        ],
    )


def _seed_owned_draft(
    *, owner, draft_id=None, fingerprint_key="k1", status=None, job_status=None, revision=1,
    job_id=None,
):
    """A pending-review draft owned by `owner`, as the worker would leave it.

    Pass `job_id` to attach the draft to a job that already exists -- the worker
    persists against the job it claimed, and a fingerprint only ever has one
    active job, so inventing a second one would not reflect any real state.
    """
    observation_id = f"obs-{uuid4().hex}"
    with db.meta_session() as session:
        session.add(models.FallbackObservation(
            id=observation_id, candidate_id=f"cand-{observation_id}",
            source_excerpt="Seven socks are laid out; which one has no pair?",
            grade_level=4, observation_kind="unsupported_shape", excluded=False,
            created_at=_now(),
        ))
        session.flush()
        job = session.get(models.GenerationJob, job_id) if job_id else None
        if job is None:
            job = models.GenerationJob(
                id=job_id or f"job-{uuid4().hex}", fingerprint_key=fingerprint_key,
                fingerprint_version=1, fingerprint_json=_fingerprint().model_dump_json(),
                trigger_observation_ids=json.dumps([observation_id]),
                status=job_status or models.JOB_SUCCEEDED, owner_session_id=owner,
                created_at=_now(), updated_at=_now(),
            )
            session.add(job)
        else:
            job.status = job_status or models.JOB_SUCCEEDED
            job.trigger_observation_ids = json.dumps([observation_id])
        session.flush()
        draft = persist_reviewable_draft(
            session, new_id=draft_id or uuid4().hex, job=job,
            candidate=_candidate_bundle(_proposal(observation_id)), now=_now(),
            revision=revision,
        )
        if status is not None:
            draft.status = status
            session.flush()
        return draft.id, job.id


# ------------------------------------------------------------- capabilities


def test_capabilities_reports_enabled_when_every_flag_is_on(client):
    assert client.get("/meta/my/capabilities").json() == {"enabled": True}


def test_capabilities_reports_disabled_without_approval(tmp_path, monkeypatch, engine):
    """Offering the button without approval would walk a teacher into a dead end."""
    _enable(monkeypatch, tmp_path, engine=engine, approval=False)
    from app.main import create_app

    client = TestClient(create_app())
    try:
        assert client.get("/meta/my/capabilities").json() == {"enabled": False}
    finally:
        get_settings.cache_clear()


def test_capabilities_reports_disabled_without_the_dynamic_classifier(tmp_path, monkeypatch, engine):
    """An approved template that /options never offers is not usable."""
    _enable(monkeypatch, tmp_path, engine=engine, dynamic_classifier=False)
    from app.main import create_app

    client = TestClient(create_app())
    try:
        assert client.get("/meta/my/capabilities").json() == {"enabled": False}
    finally:
        get_settings.cache_clear()


# ------------------------------------------------------------------ builds


def test_requesting_a_build_accepts_and_records_the_request(client):
    _start_session(client)

    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        resp = client.post("/meta/my/builds", json={"candidate_id": "c1"})

    assert resp.status_code == 202
    builds = client.get("/meta/my/builds").json()
    assert [build["candidate_id"] for build in builds] == ["c1"]


def test_requesting_a_build_without_a_session_is_rejected(client):
    resp = client.post("/meta/my/builds", json={"candidate_id": "c1"})

    assert resp.status_code == 400


def test_requesting_a_build_for_an_unknown_candidate_is_rejected(client):
    _start_session(client)

    resp = client.post("/meta/my/builds", json={"candidate_id": "nope"})

    assert resp.status_code == 404


def test_a_queued_build_reports_the_queued_stage(client):
    _start_session(client)

    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    build = client.get("/meta/my/builds").json()[0]
    assert build["stage"] == "queued"


def test_a_running_build_reports_the_building_stage(client):
    _start_session(client)
    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    with db.meta_session() as session:
        job = session.query(models.GenerationJob).one()
        job.status = models.JOB_RUNNING
        job.lease_owner = "worker-1"

    build = client.get("/meta/my/builds").json()[0]
    assert build["stage"] == "building"


def test_a_finished_build_reports_ready_with_its_draft(client):
    session_id = _start_session(client)
    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    # The worker's outcome: it claimed the queued job and left a reviewable
    # draft against it.
    with db.meta_session() as session:
        queued_job_id = session.query(models.GenerationJob).one().id
    draft_id, _ = _seed_owned_draft(
        owner=session_id,
        fingerprint_key=canonical_fingerprint_key(_fingerprint()),
        job_id=queued_job_id,
    )

    build = client.get("/meta/my/builds").json()[0]
    assert build["stage"] == "ready"
    assert build["draft_id"] == draft_id


def test_a_build_reports_a_background_failure_rather_than_waiting_forever(client):
    """A dead background task must not read as progress.

    Nothing gets enqueued if tagging exhausts its retries, so there is no job row
    to derive a stage from and the band would sit at "filed" indefinitely.
    """
    _start_session(client)

    with patch("app.meta.ingest.tag_candidate", side_effect=RuntimeError("bedrock down")):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    build = client.get("/meta/my/builds").json()[0]
    assert build["stage"] == "failed"
    assert build["error"]


def test_a_build_is_refused_when_this_session_can_already_use_a_template(client):
    session_id = _start_session(client)
    with db.meta_session() as session:
        session.add(models.TemplateVersion(
            id="tv-1", fingerprint_key=canonical_fingerprint_key(_fingerprint()),
            template_name="leftover_pair", draft_id=None, artifact_hash="sha256:x",
            status=models.TEMPLATE_VERSION_ENABLED, owner_session_id=session_id,
            created_at=_now(), updated_at=_now(),
        ))

    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    build = client.get("/meta/my/builds").json()[0]
    # Not "failed": nothing went wrong, so the band must not style it as one.
    assert build["stage"] == "already_available"
    assert "already" in build["error"].lower()


def test_a_queued_build_reports_how_long_it_has_waited(client):
    _start_session(client)
    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    build = client.get("/meta/my/builds").json()[0]
    assert build["elapsed_seconds"] >= 0


# ------------------------------------------------------------------ drafts


def test_a_draft_carries_what_a_teacher_needs_to_judge_it(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    payload = client.get(f"/meta/my/drafts/{draft_id}").json()

    assert payload["learning_objective"] == "Find the value with no partner."
    assert [beat["kind"] for beat in payload["beats"]] == ["reveal", "focus", "conclude"]
    assert payload["total_duration_seconds"] > 0
    assert payload["preview_url"] == f"/meta/my/drafts/{draft_id}/preview"
    assert payload["revision"] == 1


def test_a_draft_omits_the_machinery_a_teacher_should_not_have_to_read(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    payload = client.get(f"/meta/my/drafts/{draft_id}").json()

    for machinery in ("fixtures", "guard_document", "params_document", "validation_report",
                      "quality_report", "artifact_hash"):
        assert machinery not in payload


def test_a_draft_suggests_a_template_name_from_its_fingerprint(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    payload = client.get(f"/meta/my/drafts/{draft_id}").json()

    assert payload["suggested_template_name"] == "compare_set"


def test_another_sessions_draft_is_not_readable(client):
    _start_session(client)
    draft_id, _ = _seed_owned_draft(owner="someone-else")

    assert client.get(f"/meta/my/drafts/{draft_id}").status_code == 404


def test_an_ownerless_draft_is_not_readable_by_a_teacher(client):
    """Threshold-triggered drafts belong to the admin panel, not to a session."""
    _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=None)

    assert client.get(f"/meta/my/drafts/{draft_id}").status_code == 404


def test_a_draft_preview_is_served_to_its_owner(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    resp = client.get(f"/meta/my/drafts/{draft_id}/preview")

    assert resp.status_code == 200
    assert resp.headers["content-type"] == "image/png"


def test_another_sessions_draft_preview_is_not_served(client):
    _start_session(client)
    draft_id, _ = _seed_owned_draft(owner="someone-else")

    assert client.get(f"/meta/my/drafts/{draft_id}/preview").status_code == 404


# ------------------------------------------------------------- attempts


def test_a_refined_draft_carries_the_attempts_that_came_before_it(client):
    session_id = _start_session(client)
    first_id, job_id = _seed_owned_draft(
        owner=session_id, status=models.DRAFT_REJECTED, revision=1
    )
    with db.meta_session() as session:
        session.get(models.TemplateDraft, first_id).reviewer_feedback = "rows aren't labelled"

    second_id, _ = _seed_owned_draft(owner=session_id, revision=2)
    with db.meta_session() as session:
        second = session.get(models.TemplateDraft, second_id)
        second.parent_draft_id = first_id
        second.job_id = job_id

    payload = client.get(f"/meta/my/drafts/{second_id}").json()

    assert [attempt["revision"] for attempt in payload["attempts"]] == [1]
    assert payload["attempts"][0]["feedback"] == "rows aren't labelled"
    assert payload["attempts"][0]["preview_url"].endswith(f"/{first_id}/preview")


def test_a_first_attempt_has_no_history(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    assert client.get(f"/meta/my/drafts/{draft_id}").json()["attempts"] == []


# ------------------------------------------------------------------ verdict


def test_approving_a_draft_publishes_it_for_this_session_only(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    resp = client.post(
        f"/meta/my/drafts/{draft_id}/approve",
        json={"template_name": "leftover_pair", "math_semantics_confirmed": True},
    )

    assert resp.status_code == 200
    assert resp.json()["template_name"] == "leftover_pair"
    with db.meta_session() as session:
        version = session.query(models.TemplateVersion).one()
        assert version.owner_session_id == session_id
        assert version.status == models.TEMPLATE_VERSION_ENABLED


def test_approving_without_confirming_the_maths_is_refused(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    resp = client.post(
        f"/meta/my/drafts/{draft_id}/approve",
        json={"template_name": "leftover_pair", "math_semantics_confirmed": False},
    )

    assert resp.status_code == 422


def test_approving_another_sessions_draft_is_refused(client):
    _start_session(client)
    draft_id, _ = _seed_owned_draft(owner="someone-else")

    resp = client.post(
        f"/meta/my/drafts/{draft_id}/approve",
        json={"template_name": "leftover_pair", "math_semantics_confirmed": True},
    )

    assert resp.status_code == 404


def test_rejecting_a_draft_queues_another_attempt(client):
    session_id = _start_session(client)
    draft_id, job_id = _seed_owned_draft(owner=session_id)

    resp = client.post(
        f"/meta/my/drafts/{draft_id}/reject", json={"feedback": "the rows aren't labelled"}
    )

    assert resp.status_code == 200
    assert resp.json()["requeued"] is True
    with db.meta_session() as session:
        assert session.get(models.GenerationJob, job_id).status == models.JOB_QUEUED
        assert session.get(models.TemplateDraft, draft_id).status == models.DRAFT_REJECTED


def test_rejecting_without_a_reason_is_refused(client):
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    resp = client.post(f"/meta/my/drafts/{draft_id}/reject", json={"feedback": "   "})

    assert resp.status_code == 422


def test_rejecting_another_sessions_draft_is_refused(client):
    _start_session(client)
    draft_id, _ = _seed_owned_draft(owner="someone-else")

    resp = client.post(f"/meta/my/drafts/{draft_id}/reject", json={"feedback": "no good"})

    assert resp.status_code == 404


def test_the_teacher_routes_need_no_reviewer_token(client):
    """The admin token is exactly what this surface exists to avoid needing."""
    session_id = _start_session(client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    assert "Authorization" not in client.headers
    assert client.get(f"/meta/my/drafts/{draft_id}").status_code == 200


# ------------------------------------------------------- the capability gate
#
# GET /meta/my/capabilities already reports all four flags. The mutations have to
# honour the same gate, or a teacher can start work the deployment cannot finish:
# without codegen nothing ever drains the queue, and without the dynamic
# classifier an approved template never reaches /options.


@pytest.fixture
def codegen_disabled_client(tmp_path, monkeypatch, engine):
    _enable(monkeypatch, tmp_path, engine=engine)
    monkeypatch.setenv("META_CODEGEN_ENABLED", "0")
    get_settings.cache_clear()
    from app.main import create_app

    yield TestClient(create_app())
    get_settings.cache_clear()


def test_requesting_a_build_is_refused_when_the_generator_is_disabled(codegen_disabled_client):
    """Queuing work no worker will ever claim is worse than refusing it."""
    _start_session(codegen_disabled_client)

    resp = codegen_disabled_client.post("/meta/my/builds", json={"candidate_id": "c1"})

    assert resp.status_code == 409
    with db.meta_session() as session:
        assert session.query(models.GenerationJob).count() == 0


def test_a_build_is_not_recorded_when_the_feature_is_incomplete(codegen_disabled_client):
    _start_session(codegen_disabled_client)

    codegen_disabled_client.post("/meta/my/builds", json={"candidate_id": "c1"})

    assert codegen_disabled_client.get("/meta/my/builds").json() == []


def test_rejecting_is_refused_when_the_generator_is_disabled(codegen_disabled_client):
    """Requeueing hands work to a worker that is not running."""
    session_id = _start_session(codegen_disabled_client)
    draft_id, job_id = _seed_owned_draft(owner=session_id)

    resp = codegen_disabled_client.post(
        f"/meta/my/drafts/{draft_id}/reject", json={"feedback": "not right"}
    )

    assert resp.status_code == 409
    with db.meta_session() as session:
        assert session.get(models.GenerationJob, job_id).status == models.JOB_SUCCEEDED
        assert session.get(models.TemplateDraft, draft_id).status == models.DRAFT_PENDING_REVIEW


def test_approving_is_refused_when_the_classifier_is_disabled(tmp_path, monkeypatch, engine):
    """A published template /options can never offer is not a usable outcome."""
    _enable(monkeypatch, tmp_path, engine=engine, dynamic_classifier=False)
    from app.main import create_app

    client = TestClient(create_app())
    try:
        session_id = _start_session(client)
        draft_id, _ = _seed_owned_draft(owner=session_id)

        resp = client.post(
            f"/meta/my/drafts/{draft_id}/approve",
            json={"template_name": "leftover_pair", "math_semantics_confirmed": True},
        )

        assert resp.status_code == 409
        with db.meta_session() as session:
            assert session.query(models.TemplateVersion).count() == 0
    finally:
        get_settings.cache_clear()


def test_reading_a_draft_still_works_when_the_feature_is_incomplete(codegen_disabled_client):
    """Only mutations are gated: a teacher may still look at what exists."""
    session_id = _start_session(codegen_disabled_client)
    draft_id, _ = _seed_owned_draft(owner=session_id)

    assert codegen_disabled_client.get(f"/meta/my/drafts/{draft_id}").status_code == 200


# ------------------------------------------------- clearing a finished attempt


def test_a_teacher_can_clear_a_failed_build_and_ask_again(client):
    """No terminal state may be a dead end.

    A build that failed removes the entry button for that candidate, so without
    this the teacher has no way to try again -- which is what made the old
    global active-job refusal permanent.
    """
    _start_session(client)
    with patch("app.meta.ingest.tag_candidate", side_effect=RuntimeError("bedrock down")):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})
    assert client.get("/meta/my/builds").json()[0]["stage"] == "failed"

    resp = client.delete("/meta/my/builds/c1")

    assert resp.status_code == 204
    assert client.get("/meta/my/builds").json() == []


def test_clearing_a_build_lets_the_next_request_through(client):
    _start_session(client)
    with patch("app.meta.ingest.tag_candidate", side_effect=RuntimeError("bedrock down")):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})
    client.delete("/meta/my/builds/c1")

    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        resp = client.post("/meta/my/builds", json={"candidate_id": "c1"})

    assert resp.status_code == 202
    assert client.get("/meta/my/builds").json()[0]["stage"] == "queued"


def test_a_build_still_in_flight_cannot_be_cleared(client):
    """Clearing mid-build would orphan a queued job and hide it from its owner."""
    _start_session(client)
    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})

    resp = client.delete("/meta/my/builds/c1")

    assert resp.status_code == 409
    assert len(client.get("/meta/my/builds").json()) == 1


def test_clearing_a_build_nobody_requested_is_not_found(client):
    _start_session(client)

    assert client.delete("/meta/my/builds/c1").status_code == 404


def test_clearing_a_ready_build_is_refused(client):
    """A reviewable draft is not rubbish to sweep up; judge or reject it."""
    session_id = _start_session(client)
    with patch("app.meta.ingest.tag_candidate", return_value=_fingerprint()):
        client.post("/meta/my/builds", json={"candidate_id": "c1"})
    with db.meta_session() as session:
        queued_job_id = session.query(models.GenerationJob).one().id
    _seed_owned_draft(
        owner=session_id,
        fingerprint_key=canonical_fingerprint_key(_fingerprint()),
        job_id=queued_job_id,
    )
    assert client.get("/meta/my/builds").json()[0]["stage"] == "ready"

    assert client.delete("/meta/my/builds/c1").status_code == 409


# ------------------------------------------ cross-session ownership of status
#
# The status band derives its stage from GenerationJob rows keyed by
# fingerprint. When two sessions file the same problem shape they end up with
# distinct job rows -- the partial unique index on
# (fingerprint_key, owner_session_id) is what lets them coexist -- and neither
# session's status payload may ever surface the other's job id or draft id.
#
# These tests seed the other session's row directly rather than driving a
# second /builds request: request_template_build's real time.now() makes the
# other-session row *later* than this session's, so the plain "latest by
# fingerprint" ordering the bug used would pick it deterministically.


def _seed_foreign_job(
    *, owner: str, fingerprint_key: str, job_status: str, later_than=_now(),
) -> str:
    """A GenerationJob owned by someone else, created *after* the seeded moment.

    The bug's failure mode depends on ordering: the leaky query took the newest
    job for a fingerprint regardless of owner, so this row must post-date any
    row this session may have.
    """
    job_id = f"foreign-job-{uuid4().hex}"
    created_at = later_than + timedelta(seconds=60)
    with db.meta_session() as session:
        session.add(models.GenerationJob(
            id=job_id, fingerprint_key=fingerprint_key,
            fingerprint_version=1, fingerprint_json=_fingerprint().model_dump_json(),
            trigger_observation_ids=json.dumps([]),
            status=job_status, owner_session_id=owner,
            created_at=created_at, updated_at=created_at,
        ))
    return job_id


def _inject_template_request(session_id: str, *, candidate_id: str, fingerprint_key: str):
    """Give this session a filed request that has already been tagged.

    Bypasses the /builds background task so the test controls the timing
    relative to any foreign job it seeds afterwards.
    """
    from app.routes import store
    from app.session import TemplateRequest

    session = store.get(session_id)
    session.template_requests[candidate_id] = TemplateRequest(
        candidate_id=candidate_id,
        requested_at=_now(),
        fingerprint_key=fingerprint_key,
    )


@pytest.mark.parametrize(
    "foreign_job_status",
    [
        models.JOB_QUEUED,
        models.JOB_RUNNING,
        models.JOB_FAILED,
        models.JOB_NEEDS_MANUAL,
        models.JOB_SUCCEEDED,
    ],
)
def test_another_sessions_job_does_not_leak_into_this_sessions_status(
    client, foreign_job_status
):
    """Same fingerprint, another owner: this session must not read that stage."""
    session_id = _start_session(client)
    fp_key = canonical_fingerprint_key(_fingerprint())
    _inject_template_request(session_id, candidate_id="c1", fingerprint_key=fp_key)
    _seed_foreign_job(
        owner="someone-else", fingerprint_key=fp_key, job_status=foreign_job_status,
    )

    build = client.get("/meta/my/builds").json()[0]

    # This session has no owned job for the fingerprint, so the correct stage
    # is "filed" -- filed_but_not_yet_queued from this session's point of view.
    assert build["stage"] == "filed"
    assert build["draft_id"] is None


def test_another_sessions_ready_draft_does_not_leak_its_id(client):
    """The reported bug: session A polls, receives session B's draft id."""
    session_id = _start_session(client)
    fp_key = canonical_fingerprint_key(_fingerprint())
    _inject_template_request(session_id, candidate_id="c1", fingerprint_key=fp_key)
    # Another session has already reached a reviewable draft for the shape.
    other_draft_id, _ = _seed_owned_draft(
        owner="someone-else", fingerprint_key=fp_key, job_status=models.JOB_SUCCEEDED,
    )

    build = client.get("/meta/my/builds").json()[0]

    assert build["stage"] != "ready"
    assert build["draft_id"] != other_draft_id
    assert build["draft_id"] is None


def test_another_sessions_approved_draft_does_not_leak_its_id(client):
    """Terminal states leak too, not just in-flight ones."""
    session_id = _start_session(client)
    fp_key = canonical_fingerprint_key(_fingerprint())
    _inject_template_request(session_id, candidate_id="c1", fingerprint_key=fp_key)
    other_draft_id, _ = _seed_owned_draft(
        owner="someone-else", fingerprint_key=fp_key, status=models.DRAFT_APPROVED,
    )

    build = client.get("/meta/my/builds").json()[0]

    assert build["stage"] != "approved"
    assert build["draft_id"] != other_draft_id
    assert build["draft_id"] is None


def test_this_sessions_own_status_survives_a_newer_foreign_job(client):
    """A newer other-session job must not displace this session's own state."""
    session_id = _start_session(client)
    fp_key = canonical_fingerprint_key(_fingerprint())
    # This session's own reviewable draft, created "earlier" via _now().
    own_draft_id, own_job_id = _seed_owned_draft(
        owner=session_id, fingerprint_key=fp_key,
    )
    _inject_template_request(session_id, candidate_id="c1", fingerprint_key=fp_key)
    # Another session's newer job for the same shape. Under the bug the query
    # picked this one and either flipped the stage away from "ready" or lost
    # the draft id.
    _seed_foreign_job(
        owner="someone-else", fingerprint_key=fp_key, job_status=models.JOB_RUNNING,
    )

    build = client.get("/meta/my/builds").json()[0]

    assert build["stage"] == "ready"
    assert build["draft_id"] == own_draft_id


def test_an_ownerless_threshold_job_is_not_visible_on_the_teacher_band(client):
    """The ownerless queue is admin-scope, not this session's status."""
    session_id = _start_session(client)
    fp_key = canonical_fingerprint_key(_fingerprint())
    _inject_template_request(session_id, candidate_id="c1", fingerprint_key=fp_key)
    _seed_foreign_job(
        owner=None, fingerprint_key=fp_key, job_status=models.JOB_RUNNING,
    )

    build = client.get("/meta/my/builds").json()[0]

    assert build["stage"] == "filed"
    assert build["draft_id"] is None
