# backend/tests/test_routes.py
import io
from unittest.mock import patch

import pytest
from botocore.exceptions import NoCredentialsError
from fastapi.testclient import TestClient
from pptx import Presentation


def _pptx_bytes(slide_count: int = 1) -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for i in range(slide_count):
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i}"
        slide.placeholders[1].text = "Sarah has 4 apples and buys 3 more."
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _client():
    from app.main import create_app

    return TestClient(create_app())


def _candidate(cid="c1"):
    from app.models.candidate import Candidate

    return Candidate(
        candidate_id=cid,
        source_excerpt="Sarah has 4 apples and buys 3 more.",
        slide_index=0,
        one_line_summary="Detected: 4 + 3",
    )


def _classification():
    from app.models.scene import TemplateName
    from app.pipeline.classification import ClassificationResult, TemplateOption
    from app.templates.registry import static_ref

    return ClassificationResult(
        options=[
            TemplateOption(
                template=TemplateName.BALANCE_SCALE,
                rationale="shows the equation as a balance",
                version_id=static_ref(TemplateName.BALANCE_SCALE).version_id,
            ),
            TemplateOption(
                template=TemplateName.NUMBER_LINE,
                rationale="shows one forward jump",
                version_id=static_ref(TemplateName.NUMBER_LINE).version_id,
            ),
            TemplateOption(
                template=TemplateName.TEXT_CARD,
                rationale="always-compatible fallback",
                version_id=static_ref(TemplateName.TEXT_CARD).version_id,
            ),
        ],
        grade_level=1,
        ambiguous=False,
    )


def _upload_candidate(client):
    return _upload_candidates(client, [_candidate()])


def _upload_candidates(client, candidates):
    with patch("app.routes.discover_candidates_for_document", return_value=candidates):
        return client.post(
            "/upload",
            files={
                "file": (
                    "deck.pptx",
                    _pptx_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )


def _options_then(client):
    """Upload one candidate and cache its options; return the client."""
    with patch("app.routes.classify_candidate", return_value=_classification()):
        client.post("/options", json={"candidate_ids": ["c1"]})
    return client


def test_upload_rejects_non_pptx():
    client = _client()
    resp = client.post("/upload", files={"file": ("notes.txt", b"hello", "text/plain")})
    assert resp.status_code == 400


def test_upload_rejects_document_over_slide_cap():
    client = _client()
    with patch("app.routes.discover_candidates_for_document", return_value=[]):
        resp = client.post(
            "/upload",
            files={"file": ("big.pptx", _pptx_bytes(slide_count=51), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
    assert resp.status_code == 400


def test_upload_rejects_oversized_file():
    from app.routes import MAX_UPLOAD_BYTES

    client = _client()
    oversized = b"\x00" * (MAX_UPLOAD_BYTES + 1)
    resp = client.post(
        "/upload",
        files={"file": ("big.pptx", oversized, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert resp.status_code == 400


def test_upload_rejects_corrupt_pptx():
    client = _client()
    resp = client.post(
        "/upload",
        files={"file": ("broken.pptx", b"not a real pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )
    assert resp.status_code == 400


def test_upload_returns_candidates_and_sets_cookie():
    client = _client()
    resp = _upload_candidate(client)
    assert resp.status_code == 200
    body = resp.json()
    assert len(body["candidates"]) == 1
    assert body["candidates"][0]["candidate_id"] == "c1"
    assert "session_id" in resp.cookies


def test_upload_reports_missing_aws_credentials():
    client = _client()
    with patch(
        "app.routes.discover_candidates_for_document",
        side_effect=NoCredentialsError(),
    ):
        resp = client.post(
            "/upload",
            files={
                "file": (
                    "deck.pptx",
                    _pptx_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

    assert resp.status_code == 503
    assert resp.json() == {
        "detail": "Document analysis is unavailable because AWS credentials are not configured"
    }


def test_upload_sets_secure_cookie_when_configured():
    from app.config import Settings

    client = _client()
    with patch(
        "app.routes.get_settings",
        return_value=Settings(session_cookie_secure=True),
    ):
        resp = _upload_candidate(client)

    assert resp.status_code == 200
    assert "secure" in resp.headers["set-cookie"].lower()


def test_upload_cookie_is_not_secure_by_default():
    client = _client()
    resp = _upload_candidate(client)

    assert resp.status_code == 200
    assert "secure" not in resp.headers["set-cookie"].lower()


def test_render_without_session_is_400():
    client = _client()
    resp = client.post(
        "/render",
        json={"picks": [{"candidate_id": "c1", "template": "text_card"}]},
    )
    assert resp.status_code == 400


def test_options_returns_ranked_templates_and_caches_result():
    from app.models.scene import TemplateName
    from app.routes import store
    from app.templates.registry import static_ref

    client = _client()
    upload = _upload_candidate(client)

    with patch("app.routes.classify_candidate", return_value=_classification()) as classify:
        resp = client.post("/options", json={"candidate_ids": ["c1"]})

    assert resp.status_code == 200
    item = resp.json()["options"][0]
    static_vocab = {member.value for member in TemplateName}
    matched = {"balance_scale", "number_line", "text_card"}
    expected_rejected = sorted(static_vocab - matched)
    assert item == {
        "candidate_id": "c1",
        "grade_level": 1,
        "ambiguous": False,
        "templates": [
            {
                "template": "balance_scale",
                "version_id": static_ref(TemplateName.BALANCE_SCALE).version_id,
                "rationale": "shows the equation as a balance",
            },
            {
                "template": "number_line",
                "version_id": static_ref(TemplateName.NUMBER_LINE).version_id,
                "rationale": "shows one forward jump",
            },
            {
                "template": "text_card",
                "version_id": static_ref(TemplateName.TEXT_CARD).version_id,
                "rationale": "always-compatible fallback",
            },
        ],
        "vocabulary_size": len(static_vocab),
        "rejected": [
            {"template": name, "reason": "not_applicable"}
            for name in expected_rejected
        ],
    }
    session = store.get(upload.json()["session_id"])
    assert session.options["c1"] == _classification()
    classify.assert_called_once_with(_candidate().source_excerpt)


def test_options_response_includes_a_static_version_id():
    from app.models.scene import TemplateName
    from app.pipeline.classification import ClassificationResult, TemplateOption
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)

    number_line_ref = static_ref(TemplateName.NUMBER_LINE)
    text_card_ref = static_ref(TemplateName.TEXT_CARD)
    classification = ClassificationResult(
        options=[
            TemplateOption(
                template=TemplateName.NUMBER_LINE,
                rationale="shows one forward jump",
                version_id=number_line_ref.version_id,
            ),
            TemplateOption(
                template=TemplateName.TEXT_CARD,
                rationale="always-compatible fallback",
                version_id=text_card_ref.version_id,
            ),
        ],
        grade_level=1,
        ambiguous=False,
    )
    with patch("app.routes.classify_candidate", return_value=classification):
        resp = client.post("/options", json={"candidate_ids": ["c1"]})

    assert resp.status_code == 200
    number_line = next(
        t for t in resp.json()["options"][0]["templates"] if t["template"] == "number_line"
    )
    assert number_line["version_id"] == number_line_ref.version_id


def test_options_marks_all_structural_templates_low_confidence_when_ambiguous():
    from app.models.scene import TemplateName
    from app.pipeline.classification import ClassificationResult, TemplateOption
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)

    # Ambiguity forces classify_candidate to strip structural options and keep
    # only the text_card fallback, so every non-text_card template in the vocab
    # ends up rejected with reason=low_confidence.
    text_card_ref = static_ref(TemplateName.TEXT_CARD)
    classification = ClassificationResult(
        options=[
            TemplateOption(
                template=TemplateName.TEXT_CARD,
                rationale="always-compatible fallback",
                version_id=text_card_ref.version_id,
            ),
        ],
        grade_level=2,
        ambiguous=True,
    )
    with patch("app.routes.classify_candidate", return_value=classification):
        resp = client.post("/options", json={"candidate_ids": ["c1"]})

    assert resp.status_code == 200
    item = resp.json()["options"][0]
    static_vocab = {member.value for member in TemplateName}
    assert item["vocabulary_size"] == len(static_vocab)
    expected = sorted(static_vocab - {"text_card"})
    assert item["rejected"] == [
        {"template": name, "reason": "low_confidence"} for name in expected
    ]


def test_options_includes_dynamic_templates_in_vocabulary_when_enabled():
    from app.config import get_settings
    from app.meta.dynamic_templates import DynamicSnapshotEntry, EnabledSnapshot

    client = _client()
    _upload_candidate(client)

    entries = {
        name: DynamicSnapshotEntry(
            version_id=f"{name}-v1",
            artifact_hash=f"{name}-hash",
            classifier_bullet=f"- {name}: fake",
        )
        for name in ("custom_bar", "custom_grid")
    }
    snapshot = EnabledSnapshot(_entries=entries)

    settings = get_settings()
    settings.meta_dynamic_classifier_enabled = True
    try:
        with patch("app.routes.meta_session") as mock_meta_session, patch(
            "app.routes.load_enabled_snapshot", return_value=snapshot
        ), patch(
            "app.routes.classify_candidate", return_value=_classification()
        ):
            mock_meta_session.return_value.__enter__.return_value = object()
            resp = client.post("/options", json={"candidate_ids": ["c1"]})
    finally:
        settings.meta_dynamic_classifier_enabled = False

    assert resp.status_code == 200
    item = resp.json()["options"][0]
    from app.models.scene import TemplateName

    expected_vocab = {m.value for m in TemplateName} | set(entries)
    assert item["vocabulary_size"] == len(expected_vocab)
    rejected_names = {r["template"] for r in item["rejected"]}
    # The dynamic templates the model did not pick appear as rejected.
    assert set(entries) <= rejected_names


def test_options_passes_no_session_when_dynamic_classifier_flag_is_off():
    from unittest.mock import patch

    client = _client()
    _upload_candidate(client)

    with patch("app.routes.classify_candidate") as mock_classify:
        from app.pipeline.classification import ClassificationResult

        mock_classify.return_value = ClassificationResult(options=[], grade_level=1, ambiguous=False)
        client.post("/options", json={"candidate_ids": ["c1"]})

    assert mock_classify.call_args.kwargs.get("session") is None


def test_options_loads_one_snapshot_for_the_whole_batch_when_flag_enabled():
    from app.config import get_settings
    from app.meta.dynamic_templates import EnabledSnapshot

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), _candidate("c3")])

    settings = get_settings()
    settings.meta_dynamic_classifier_enabled = True
    try:
        with patch("app.routes.meta_session") as mock_meta_session, patch(
            "app.routes.load_enabled_snapshot", return_value=EnabledSnapshot(_entries={})
        ) as mock_load_snapshot, patch(
            "app.routes.classify_candidate", return_value=_classification()
        ) as mock_classify:
            mock_meta_session.return_value.__enter__.return_value = object()
            resp = client.post(
                "/options", json={"candidate_ids": ["c1", "c2", "c3"]}
            )
    finally:
        settings.meta_dynamic_classifier_enabled = False

    assert resp.status_code == 200
    assert mock_meta_session.call_count == 1
    assert mock_load_snapshot.call_count == 1
    assert mock_classify.call_count == 3
    for call in mock_classify.call_args_list:
        assert call.kwargs.get("snapshot") is mock_load_snapshot.return_value


def test_storyboard_resolves_a_dynamic_template_pick():
    from unittest.mock import MagicMock, patch

    from app.models.scene import Scene, TemplateRef
    from app.pipeline.classification import ClassificationResult, TemplateOption

    client = _client()
    _upload_candidate(client)

    classification = ClassificationResult(
        options=[
            TemplateOption(template="decimal_comparison_grid", rationale="fits", version_id="v1"),
            TemplateOption(
                template="text_card",
                rationale="always-compatible fallback",
                version_id="static:text_card:1",
            ),
        ],
        grade_level=4,
        ambiguous=False,
    )
    dynamic_ref = TemplateRef(name="decimal_comparison_grid", version_id="v1", artifact_hash="sha256:x")

    fake_params_cls = MagicMock()
    fake_params_cls.model_json_schema.return_value = {}

    with patch("app.routes.classify_candidate", return_value=classification), patch(
        "app.routes.resolve_dynamic_ref", return_value=dynamic_ref
    ) as mock_resolve_dynamic, patch("app.routes.assemble_scene") as mock_assemble, patch(
        "app.meta.dynamic_templates.get_dynamic_template",
        return_value=(MagicMock(), fake_params_cls),
    ):
        mock_assemble.return_value = Scene(
            scene_id="s1",
            candidate_id="c1",
            template=dynamic_ref,
            grade_level=4,
            params={},
            status="pending_review",
        )
        client.post("/options", json={"candidate_ids": ["c1"]})
        resp = client.post(
            "/storyboard",
            json={"picks": [{"candidate_id": "c1", "template": "decimal_comparison_grid"}]},
        )

    assert resp.status_code == 200
    mock_resolve_dynamic.assert_called_once()


def test_options_unknown_candidate_is_404():
    client = _client()
    _upload_candidate(client)

    resp = client.post("/options", json={"candidate_ids": ["does-not-exist"]})

    assert resp.status_code == 404


def test_options_without_session_is_400():
    client = _client()

    resp = client.post("/options", json={"candidate_ids": ["c1"]})

    assert resp.status_code == 400


def test_options_rejects_duplicate_candidates_before_classification():
    client = _client()
    _upload_candidate(client)

    with patch("app.routes.classify_candidate", return_value=_classification()) as classify:
        resp = client.post("/options", json={"candidate_ids": ["c1", "c1"]})

    assert resp.status_code == 400
    assert "duplicate" in resp.json()["detail"].lower()
    classify.assert_not_called()


def test_options_rejects_more_than_50_candidates_before_classification():
    client = _client()
    _upload_candidate(client)

    with patch("app.routes.classify_candidate") as classify:
        resp = client.post(
            "/options",
            json={"candidate_ids": [f"c{i}" for i in range(51)]},
        )

    assert resp.status_code == 422
    classify.assert_not_called()


def test_unknown_clip_id_is_404():
    client = _client()
    resp = client.get("/clips/nope")
    assert resp.status_code == 404


def test_storyboard_builds_scenes_with_schema_and_thumbnail_url(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)
    _options_then(client)

    thumb = tmp_path / "t.png"
    thumb.write_bytes(b"png")
    fake = Scene(
        scene_id="s1",
        candidate_id="c1",
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        status="pending_review",
        thumbnail_path=thumb,
    )

    with patch("app.routes.assemble_scene", return_value=fake):
        resp = client.post(
            "/storyboard",
            json={"picks": [{"candidate_id": "c1", "template": "number_line"}]},
        )

    assert resp.status_code == 200
    scene = resp.json()["scenes"][0]
    assert scene["scene_id"] == "s1"
    assert scene["status"] == "pending_review"
    assert scene["thumbnail_url"].startswith("/thumbnails/")
    assert scene["source_excerpt"]
    assert scene["detected_summary"] == "Detected: 4 + 3"
    assert scene["params_schema"]["properties"]["start"]["type"] == "integer"
    # Compile artifact — surfaces the pure params → scene-program step.
    assert isinstance(scene["scene_program_hash"], str)
    assert len(scene["scene_program_hash"]) == 64
    assert scene["program_size"] > 0
    assert scene["compile_ms"] is not None
    assert scene["scene_program"]["params"] == scene["params"]
    assert scene["scene_program"]["template"]["name"] == "number_line"
    # Named gate list on SceneOut: audience-facing surfaces render this so N
    # named gates are visible next to the deck stamps, not just generic ones.
    gate_names = [g["name"] for g in scene["gates"]]
    assert gate_names == [
        "Values extracted",
        "Schema check",
        "Semantic check",
        "Compiled deterministically",
        "Preview rendered",
    ]
    categories = {g["category"] for g in scene["gates"]}
    assert categories <= {"Fixture", "Anchor alignment", "Rendered output"}
    assert all(g["status"] == "passed" for g in scene["gates"])


def test_storyboard_exposes_stated_answer_and_mismatch(tmp_path):
    from fractions import Fraction

    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)
    _options_then(client)

    thumb = tmp_path / "t.png"
    thumb.write_bytes(b"png")
    fake = Scene(
        scene_id="s1",
        candidate_id="c1",
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        status="pending_review",
        thumbnail_path=thumb,
        stated_answer=Fraction(9),
        stated_answer_source="= 9",
    )

    with patch("app.routes.assemble_scene", return_value=fake):
        resp = client.post(
            "/storyboard",
            json={"picks": [{"candidate_id": "c1", "template": "number_line"}]},
        )

    assert resp.status_code == 200
    scene = resp.json()["scenes"][0]
    assert scene["stated_answer"] == "9"
    assert scene["stated_answer_source"] == "= 9"
    assert scene["mismatch"] == {"stated": "9", "computed": "7"}
    assert scene["mismatch_acknowledged"] is False


def test_storyboard_does_not_break_with_meta_flag_off(tmp_path, monkeypatch):
    # meta_templates_enabled defaults to False, so the storyboard route must behave
    # exactly as before — record_unsupported_shape returns immediately. We
    # monkeypatch it here only to observe that the wiring calls it; the real
    # flag-off short-circuit is covered separately in tests/meta/test_ingest.py.
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    calls = []
    import app.routes as routes

    monkeypatch.setattr(routes, "record_unsupported_shape", lambda **kw: calls.append(kw))

    client = _client()
    _upload_candidate(client)
    _options_then(client)

    thumb = tmp_path / "t.png"
    thumb.write_bytes(b"png")
    fake = Scene(
        scene_id="s1",
        candidate_id="c1",
        template=static_ref(TemplateName.TEXT_CARD),
        grade_level=1,
        params={"headline": "x", "lines": ["y"]},
        status="pending_review",
        thumbnail_path=thumb,
    )

    with patch("app.routes.assemble_scene", return_value=fake):
        resp = client.post(
            "/storyboard",
            json={"picks": [{"candidate_id": "c1", "template": "text_card"}]},
        )

    assert resp.status_code == 200
    assert len(calls) == 1
    assert calls[0]["candidate_id"] == "c1"
    assert calls[0]["picked_template"] == TemplateName.TEXT_CARD
    assert calls[0]["scene_status"] == "pending_review"


def test_thumbnail_endpoint_serves_png(tmp_path):
    from app.routes import store

    client = _client()
    png = tmp_path / "t.png"
    png.write_bytes(b"\x89PNG\r\n")
    thumb_id = store.register_thumbnail(png)

    resp = client.get(f"/thumbnails/{thumb_id}")
    assert resp.status_code == 200
    assert resp.headers["content-type"] == "image/png"


def test_thumbnail_unknown_id_is_404():
    client = _client()
    assert client.get("/thumbnails/nope").status_code == 404


def test_storyboard_rejects_pick_before_options_cached():
    client = _client()
    _upload_candidate(client)
    resp = client.post(
        "/storyboard",
        json={"picks": [{"candidate_id": "c1", "template": "number_line"}]},
    )
    assert resp.status_code == 400


def test_storyboard_without_session_is_400():
    client = _client()
    resp = client.post(
        "/storyboard",
        json={"picks": [{"candidate_id": "c1", "template": "number_line"}]},
    )
    assert resp.status_code == 400


def test_storyboard_stale_cached_version_id_is_409():
    from app.models.scene import TemplateName, TemplateVersionMismatchError
    from app.pipeline.classification import ClassificationResult, TemplateOption

    client = _client()
    _upload_candidate(client)

    stale_classification = ClassificationResult(
        options=[
            TemplateOption(
                template=TemplateName.NUMBER_LINE,
                rationale="shows one forward jump",
                version_id="stale-version",
            ),
        ],
        grade_level=1,
        ambiguous=False,
    )
    with patch("app.routes.classify_candidate", return_value=stale_classification):
        client.post("/options", json={"candidate_ids": ["c1"]})

    with patch(
        "app.routes.resolve_static_ref",
        side_effect=TemplateVersionMismatchError("stale contract"),
    ):
        resp = client.post(
            "/storyboard",
            json={"picks": [{"candidate_id": "c1", "template": "number_line"}]},
        )

    assert resp.status_code == 409


def _seed_scene(client, scene, template=None):
    """Attach `scene` to the client's current session (in the module-level store)."""
    from app.models.scene import TemplateName
    from app.routes import store
    from app.templates.registry import static_ref

    session_id = client.cookies.get("session_id")
    session = store.get(session_id)
    session.scenes[scene.scene_id] = scene
    session.scene_order.append(scene.scene_id)
    if template is not None:
        session.scene_requested_template[scene.scene_id] = static_ref(TemplateName(template))
    return session


def _number_line_scene(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    thumb = tmp_path / "t.png"
    thumb.write_bytes(b"png")
    return Scene(
        scene_id="s1",
        candidate_id="c1",
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        status="pending_review",
        thumbnail_path=thumb,
    )


def test_render_renders_only_approved_from_stored_params(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path)
    approved = approved.model_copy(update={"status": "approved"})
    _seed_scene(client, approved)

    def fake_render(template, params, out):
        out.write_bytes(b"mp4")
        return out

    # Bedrock extraction must NOT be called at render time.
    with patch("app.routes.render_scene_to_mp4", side_effect=fake_render), patch(
        "app.pipeline.process_scene.extract_params"
    ) as extract:
        resp = client.post("/render")

    assert resp.status_code == 200
    clips = resp.json()["clips"]
    assert len(clips) == 1
    assert clips[0]["scene_id"] == "s1"
    assert clips[0]["clip_url"].startswith("/clips/")
    extract.assert_not_called()


def test_render_returns_manual_scene_results(tmp_path):
    from app.main import create_app
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = TestClient(create_app(), raise_server_exceptions=False)
    _upload_candidate(client)
    manual = Scene(
        scene_id="manual-1",
        manual_source_text="Show 3 + 4 on a number line.",
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"start": 3, "steps": [{"operation": "add", "amount": 4}]},
        status="approved",
    )
    _seed_scene(client, manual)

    def fake_render(template, params, out):
        out.write_bytes(b"mp4")
        return out

    with patch("app.routes.render_scene_to_mp4", side_effect=fake_render):
        resp = client.post("/render")

    assert resp.status_code == 200
    clip = resp.json()["clips"][0]
    gates = clip.pop("gates")
    assert clip == {
        "scene_id": "manual-1",
        "candidate_id": None,
        "candidate_ids": None,
        "status": "approved",
        "clip_url": clip["clip_url"],
        "fallback_reason": None,
    }
    assert [g["name"] for g in gates][-1] == "Full render"
    assert next(g["status"] for g in gates if g["name"] == "Full render") == "passed"


def test_render_skips_rejected_scenes(tmp_path):
    client = _client()
    _upload_candidate(client)
    rejected = _number_line_scene(tmp_path).model_copy(update={"status": "rejected"})
    _seed_scene(client, rejected)
    resp = client.post("/render")
    assert resp.status_code == 400  # nothing approved


def test_render_one_failure_does_not_sink_batch(tmp_path):
    from app.models.scene import Scene, TemplateName

    client = _client()
    _upload_candidate(client)
    good = _number_line_scene(tmp_path).model_copy(
        update={"scene_id": "sg", "status": "approved"}
    )
    bad = _number_line_scene(tmp_path).model_copy(
        update={"scene_id": "sb", "status": "approved"}
    )
    _seed_scene(client, good)
    _seed_scene(client, bad)

    calls = {"n": 0}

    def render_side_effect(template, params, out):
        calls["n"] += 1
        if calls["n"] == 2:
            raise RuntimeError("boom")
        out.write_bytes(b"mp4")
        return out

    with patch("app.routes.render_scene_to_mp4", side_effect=render_side_effect):
        resp = client.post("/render")

    assert resp.status_code == 200
    statuses = {c["status"] for c in resp.json()["clips"]}
    assert "error" in statuses
    assert len(resp.json()["clips"]) == 2


def test_render_stored_param_validation_failure_does_not_sink_batch(tmp_path):
    good = _number_line_scene(tmp_path).model_copy(
        update={"scene_id": "sg", "status": "approved"}
    )
    bad = _number_line_scene(tmp_path).model_copy(
        update={
            "scene_id": "sb",
            "status": "approved",
            # Guard-invalid: running total goes negative (1 - 5 = -4).
            "params": {"start": 1, "steps": [{"operation": "subtract", "amount": 5}]},
        }
    )

    client = _client()
    _upload_candidate(client)
    _seed_scene(client, good)
    _seed_scene(client, bad)

    def fake_render(template, params, out):
        out.write_bytes(b"mp4")
        return out

    with patch("app.routes.render_scene_to_mp4", side_effect=fake_render):
        resp = client.post("/render")

    assert resp.status_code == 200
    clips = resp.json()["clips"]
    assert len(clips) == 2
    # Scenes render in scene_order (good seeded first, bad second).
    statuses = [c["status"] for c in clips]
    assert statuses == ["approved", "error"]
    assert clips[1]["clip_url"] is None


def test_patch_valid_params_re_renders_thumbnail(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    with patch("app.routes.render_scene_thumbnail") as thumb:
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": 10, "steps": [{"operation": "subtract", "amount": 2}]}},
        )

    assert resp.status_code == 200
    body = resp.json()
    assert body["params"]["start"] == 10
    assert body["status"] == "pending_review"
    thumb.assert_called_once()


def test_patch_invalid_params_returns_422_and_keeps_scene(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    # start=1 then subtract 5 -> running total goes negative -> guard rejects.
    with patch("app.routes.render_scene_thumbnail") as thumb:
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": 1, "steps": [{"operation": "subtract", "amount": 5}]}},
        )

    assert resp.status_code == 422
    errors = resp.json()["detail"]["errors"]
    assert errors
    # The @model_validator guard rejects, so pydantic tags this value_error —
    # UI uses `category` to route to the "Semantic check" stamp and `rule`
    # as the per-rule identifier (rather than a shared value_error label).
    assert errors[0]["type"] == "value_error"
    assert errors[0]["category"] == "semantic"
    assert errors[0]["rule"] and errors[0]["rule"] != "value_error"
    thumb.assert_not_called()


def test_patch_failed_thumbnail_regen_leaves_old_approval_intact(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    with patch("app.routes.render_scene_thumbnail", side_effect=RuntimeError("boom")):
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": 10, "steps": [{"operation": "subtract", "amount": 2}]}},
        )

    assert resp.status_code == 500

    from app.routes import store

    session = store.get(client.cookies["session_id"])
    scene = session.scenes["s1"]
    assert scene.status == "approved"
    assert scene.params["start"] == 4  # original params, edit never committed
    assert scene.revision == 0


def test_patch_grade_sets_overridden(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    resp = client.patch("/storyboard/s1", json={"grade_level": 5})
    assert resp.status_code == 200
    assert resp.json()["grade_level"] == 5
    assert resp.json()["grade_overridden"] is True


def test_patch_wrong_param_type_returns_schema_type_error(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    with patch("app.routes.render_scene_thumbnail") as thumb:
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": "not-an-int", "steps": []}},
        )

    assert resp.status_code == 422
    errors = resp.json()["detail"]["errors"]
    assert errors
    # Wrong scalar type is a schema-level failure; UI routes non-value_error /
    # non-assertion_error entries to the "Schema check" stamp.
    types = {e["type"] for e in errors}
    assert any(t not in {"value_error", "assertion_error"} for t in types)
    thumb.assert_not_called()


def test_patch_out_of_range_grade_returns_field_errors_shape(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    with patch("app.routes.render_scene_thumbnail") as thumb:
        resp = client.patch("/storyboard/s1", json={"grade_level": 100})

    assert resp.status_code == 422
    errors = resp.json()["detail"]["errors"]
    assert errors
    assert "loc" in errors[0]
    assert "msg" in errors[0]
    assert "type" in errors[0]
    # Grade range is a range/schema check, not a cross-field rule.
    assert errors[0]["category"] == "schema"
    assert errors[0]["rule"] == "grade_range"
    thumb.assert_not_called()


def test_patch_unknown_scene_is_404():
    client = _client()
    _upload_candidate(client)
    resp = client.patch("/storyboard/nope", json={"grade_level": 3})
    assert resp.status_code == 404


def test_retry_reextracts_same_template_and_keeps_scene_id(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path), template="number_line")

    fresh = Scene(
        scene_id="ignored-new-id",
        candidate_id="c1",
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        status="pending_review",
        thumbnail_path=(tmp_path / "t.png"),
    )

    with patch("app.routes.assemble_scene", return_value=fresh) as assemble:
        resp = client.post("/storyboard/s1/retry")

    assert resp.status_code == 200
    assert resp.json()["scene_id"] == "s1"  # replaced in place
    # retried on the originally-picked template
    assert assemble.call_args.kwargs["template"] == static_ref(TemplateName.NUMBER_LINE)


def test_retry_unknown_scene_is_404():
    client = _client()
    _upload_candidate(client)
    assert client.post("/storyboard/nope/retry").status_code == 404


def test_approve_sets_status(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))
    resp = client.post("/storyboard/s1/approve")
    assert resp.status_code == 200
    assert resp.json()["status"] == "approved"


def test_reject_sets_status(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))
    resp = client.post("/storyboard/s1/reject")
    assert resp.status_code == 200
    assert resp.json()["status"] == "rejected"


def test_approve_fallback_scene_keeps_reason(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidate(client)
    fallback = Scene(
        scene_id="s2",
        candidate_id="c1",
        template=static_ref(TemplateName.TEXT_CARD),
        grade_level=1,
        params={"headline": "x", "lines": ["y"]},
        status="fallback",
        fallback_reason="did not fit the chosen template",
        thumbnail_path=(tmp_path / "t.png"),
    )
    (tmp_path / "t.png").write_bytes(b"png")
    _seed_scene(client, fallback)

    resp = client.post("/storyboard/s2/approve")
    assert resp.status_code == 200
    assert resp.json()["status"] == "approved"
    assert resp.json()["fallback_reason"] == "did not fit the chosen template"


def test_approve_unknown_scene_is_404():
    client = _client()
    _upload_candidate(client)
    assert client.post("/storyboard/nope/approve").status_code == 404


def test_approve_chained_scene_returns_candidate_ids_and_joined_text(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    thumb = tmp_path / "t.png"
    thumb.write_bytes(b"png")
    chained = Scene(
        scene_id="s1",
        candidate_ids=["c1", "c2"],
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"items": [
            {"start": 4, "steps": [{"operation": "add", "amount": 3}]},
            {"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        ]},
        status="pending_review",
        thumbnail_path=thumb,
    )
    _seed_scene(client, chained)

    resp = client.post("/storyboard/s1/approve")

    assert resp.status_code == 200
    body = resp.json()
    assert body["candidate_ids"] == ["c1", "c2"]
    assert body["candidate_id"] is None
    assert body["source_excerpt"] == (
        "Sarah has 4 apples and buys 3 more. / Sarah has 4 apples and buys 3 more."
    )
    assert body["detected_summary"] == "Detected: 4 + 3 / Detected: 4 + 3"
    assert body["params_schema"]["properties"]["items"]["type"] == "array"


def _mismatched_scene(tmp_path):
    """Number-line scene whose params compute 8 but which states an answer of 9."""
    from fractions import Fraction

    scene = _number_line_scene(tmp_path).model_copy(
        update={
            "params": {"start": 4, "steps": [{"operation": "add", "amount": 4}]},
            "stated_answer": Fraction(9),
            "stated_answer_source": "= 9",
        }
    )
    return scene


def _matching_scene(tmp_path):
    """Number-line scene whose stated answer agrees with its computed answer (7)."""
    from fractions import Fraction

    return _number_line_scene(tmp_path).model_copy(
        update={"stated_answer": Fraction(7), "stated_answer_source": "= 7"}
    )


def test_approve_blocked_by_unacked_mismatch(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _mismatched_scene(tmp_path))

    resp = client.post("/storyboard/s1/approve")

    assert resp.status_code == 409
    body = resp.json()["detail"]
    assert body["error"] == "stated_answer_mismatch"
    assert body["stated"] == "9"
    assert body["computed"] == "8"


def test_approve_succeeds_when_mismatch_acknowledged(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _mismatched_scene(tmp_path))

    ack = client.post("/storyboard/s1/acknowledge-mismatch")
    assert ack.status_code == 200

    approve = client.post("/storyboard/s1/approve")
    assert approve.status_code == 200
    assert approve.json()["status"] == "approved"


def test_approve_succeeds_when_answers_match(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _matching_scene(tmp_path))

    resp = client.post("/storyboard/s1/approve")
    assert resp.status_code == 200


def test_approve_succeeds_when_no_stated_answer(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    resp = client.post("/storyboard/s1/approve")
    assert resp.status_code == 200


def test_acknowledge_mismatch_without_mismatch_returns_409(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _matching_scene(tmp_path))

    resp = client.post("/storyboard/s1/acknowledge-mismatch")
    assert resp.status_code == 409


def test_acknowledge_mismatch_unknown_scene_is_404():
    client = _client()
    _upload_candidate(client)
    assert client.post("/storyboard/nope/acknowledge-mismatch").status_code == 404


def test_edit_resets_mismatch_acknowledged(tmp_path):
    from app.routes import store

    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _mismatched_scene(tmp_path))

    ack = client.post("/storyboard/s1/acknowledge-mismatch")
    assert ack.status_code == 200
    session = store.get(client.cookies["session_id"])
    assert session.scenes["s1"].mismatch_acknowledged is True

    # Edit params: content changed, so the ack resets regardless of the new mismatch state.
    with patch("app.routes.render_scene_thumbnail"):
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": 4, "steps": [{"operation": "add", "amount": 5}]}},
        )

    assert resp.status_code == 200
    assert resp.json()["mismatch_acknowledged"] is False
    assert session.scenes["s1"].mismatch_acknowledged is False


def test_edit_noop_patch_does_not_reset_mismatch_acknowledged(tmp_path):
    from app.routes import store

    client = _client()
    _upload_candidate(client)
    scene = _mismatched_scene(tmp_path)
    _seed_scene(client, scene)

    ack = client.post("/storyboard/s1/acknowledge-mismatch")
    assert ack.status_code == 200

    # No-op PATCH: resend the same grade_level, no params change.
    resp = client.patch("/storyboard/s1", json={"grade_level": scene.grade_level})

    assert resp.status_code == 200
    assert resp.json()["mismatch_acknowledged"] is True
    session = store.get(client.cookies["session_id"])
    assert session.scenes["s1"].mismatch_acknowledged is True


def test_chained_scene_mismatch_uses_last_item_compute_answer(tmp_path):
    """compute_answer_for routes candidate_ids scenes through get_chained_template;
    the mismatch must reflect the chained params' last item, not a solo template."""
    from fractions import Fraction

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    # Last item computes 5 - 1 = 4, but the source states the answer is 9.
    chained = _chained_number_line_scene().model_copy(
        update={"stated_answer": Fraction(9), "stated_answer_source": "= 9"}
    )
    _seed_scene(client, chained)

    resp = client.patch("/storyboard/s1", json={"grade_level": chained.grade_level})

    assert resp.status_code == 200
    assert resp.json()["mismatch"] == {"stated": "9", "computed": "4"}


def test_combine_rejects_scene_with_unacked_mismatch(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    scene1 = _mismatched_scene(tmp_path)
    scene2 = _mismatched_scene(tmp_path).model_copy(
        update={"scene_id": "s2", "candidate_id": "c2"}
    )
    _seed_scene(client, scene1)
    _seed_scene(client, scene2)

    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})

    assert resp.status_code == 409
    body = resp.json()["detail"]
    assert body["error"] == "stated_answer_mismatch_in_chain"
    assert body["scene_id"] == "s1"


def test_combine_succeeds_when_mismatch_acknowledged(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    scene1 = _mismatched_scene(tmp_path)
    scene2 = _mismatched_scene(tmp_path).model_copy(
        update={"scene_id": "s2", "candidate_id": "c2"}
    )
    _seed_scene(client, scene1)
    _seed_scene(client, scene2)

    assert client.post("/storyboard/s1/acknowledge-mismatch").status_code == 200
    assert client.post("/storyboard/s2/acknowledge-mismatch").status_code == 200

    with patch("app.routes.render_chained_scene_thumbnail"):
        resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})

    assert resp.status_code == 200


def test_patch_params_resets_approved_scene_to_pending_review(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    with patch("app.routes.render_scene_thumbnail"):
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"start": 10, "steps": [{"operation": "subtract", "amount": 2}]}},
        )

    assert resp.status_code == 200
    assert resp.json()["status"] == "pending_review"


def test_patch_grade_resets_approved_scene_to_pending_review(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    resp = client.patch("/storyboard/s1", json={"grade_level": 5})

    assert resp.status_code == 200
    assert resp.json()["status"] == "pending_review"


def test_patch_empty_body_does_not_revoke_approval(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    resp = client.patch("/storyboard/s1", json={})

    assert resp.status_code == 200
    assert resp.json()["status"] == "approved"

    with patch("app.routes.render_scene_to_mp4", side_effect=lambda t, p, out: out.write_bytes(b"mp4")):
        render_resp = client.post("/render")
    assert render_resp.status_code == 200


def test_patch_same_grade_level_does_not_revoke_approval(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    resp = client.patch("/storyboard/s1", json={"grade_level": approved.grade_level})

    assert resp.status_code == 200
    assert resp.json()["status"] == "approved"


def test_patch_same_params_does_not_revoke_approval(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    with patch("app.routes.render_scene_thumbnail"):
        resp = client.patch("/storyboard/s1", json={"params": approved.params})

    assert resp.status_code == 200
    assert resp.json()["status"] == "approved"


def test_patch_resets_rejected_scene_to_pending_review(tmp_path):
    client = _client()
    _upload_candidate(client)
    rejected = _number_line_scene(tmp_path).model_copy(update={"status": "rejected"})
    _seed_scene(client, rejected)

    resp = client.patch("/storyboard/s1", json={"grade_level": 5})

    assert resp.status_code == 200
    assert resp.json()["status"] == "pending_review"


def test_edited_approved_scene_cannot_render_until_reapproved(tmp_path):
    client = _client()
    _upload_candidate(client)
    approved = _number_line_scene(tmp_path).model_copy(
        update={"status": "approved", "approved_revision": 0}
    )
    _seed_scene(client, approved)

    with patch("app.routes.render_scene_thumbnail"):
        client.patch(
            "/storyboard/s1",
            json={"params": {"start": 10, "steps": [{"operation": "subtract", "amount": 2}]}},
        )

    resp = client.post("/render")
    assert resp.status_code == 400  # nothing approved


def test_approve_records_revision_edit_after_approve_does_not_authorize_render(tmp_path):
    """An approval for an old revision cannot authorize a newer, unapproved revision.

    Simulates a hypothetical mutation path that changes render-affecting scene
    content without resetting status (defense in depth beyond the PATCH-side
    status reset already covered above).
    """
    from app.routes import store

    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    approve_resp = client.post("/storyboard/s1/approve")
    assert approve_resp.status_code == 200

    session = store.get(client.cookies["session_id"])
    stale = session.scenes["s1"]
    assert stale.approved_revision == stale.revision
    session.scenes["s1"] = stale.model_copy(
        update={"params": {"start": 99, "steps": []}, "revision": stale.revision + 1}
    )

    def fake_render(template, params, out):
        out.write_bytes(b"mp4")
        return out

    with patch("app.routes.render_scene_to_mp4", side_effect=fake_render):
        resp = client.post("/render")

    assert resp.status_code == 400  # approval no longer matches current revision


def test_concurrent_edit_and_approve_conflicts(tmp_path, monkeypatch):
    """A late edit and a simultaneous approve must not silently pick a winner."""
    import threading

    from app import routes

    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    started = threading.Event()
    proceed = threading.Event()
    original_cas = routes._write_scene_cas

    def delayed_cas(session, scene_id, base_revision, updates):
        started.set()
        assert proceed.wait(timeout=5)
        return original_cas(session, scene_id, base_revision, updates)

    monkeypatch.setattr(routes, "_write_scene_cas", delayed_cas)

    results = {}

    def do_approve():
        results["approve"] = client.post("/storyboard/s1/approve")

    thread = threading.Thread(target=do_approve)
    thread.start()
    assert started.wait(timeout=5)

    monkeypatch.setattr(routes, "_write_scene_cas", original_cas)
    edit_resp = client.patch("/storyboard/s1", json={"grade_level": 6})
    assert edit_resp.status_code == 200

    proceed.set()
    thread.join(timeout=5)

    assert results["approve"].status_code == 409


def test_chain_combines_two_scenes_into_one(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    scene1 = Scene(
        scene_id="s1", candidate_id="c1", template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1, params={"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        status="pending_review",
    )
    scene2 = Scene(
        scene_id="s2", candidate_id="c2", template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1, params={"start": 5, "steps": [{"operation": "subtract", "amount": 1}]},
        status="pending_review",
    )
    _seed_scene(client, scene1)
    _seed_scene(client, scene2)

    with patch("app.routes.render_chained_scene_thumbnail") as thumb:
        resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})

    assert resp.status_code == 200
    body = resp.json()
    assert body["candidate_ids"] == ["c1", "c2"]
    assert body["status"] == "pending_review"
    assert body["template"] == "number_line"
    thumb.assert_called_once()

    from app.routes import store
    session = store.get(client.cookies.get("session_id"))
    assert session.scene_order == [body["scene_id"]]
    assert session.scene_chain_members[body["scene_id"]] == ["s1", "s2"]


def test_chain_rejects_fewer_than_two_ids():
    client = _client()
    _upload_candidate(client)
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1"]})
    assert resp.status_code == 400
    assert resp.json()["detail"] == "A chain must contain between 2 and 4 scenes"


def test_chain_rejects_more_than_four_ids():
    client = _client()
    _upload_candidate(client)
    resp = client.post(
        "/storyboard/chain", json={"scene_ids": ["s1", "s2", "s3", "s4", "s5"]}
    )
    assert resp.status_code == 400
    assert resp.json()["detail"] == "A chain must contain between 2 and 4 scenes"


def test_chain_rejects_duplicate_ids(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s1"]})
    assert resp.status_code == 400
    assert "duplicate" in resp.json()["detail"].lower()


def test_chain_rejects_unknown_scene_id(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "does-not-exist"]})
    assert resp.status_code == 400


def test_chain_rejects_non_pending_review_scene(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    approved = _number_line_scene(tmp_path).model_copy(update={"status": "approved"})
    pending = _number_line_scene(tmp_path).model_copy(
        update={"scene_id": "s2", "candidate_id": "c2"}
    )
    _seed_scene(client, approved)
    _seed_scene(client, pending)
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})
    assert resp.status_code == 400


def test_chain_rejects_already_chained_scene(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    already_chained = Scene(
        scene_id="s1", candidate_ids=["ca", "cb"], template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"items": [
            {"start": 4, "steps": [{"operation": "add", "amount": 3}]},
            {"start": 4, "steps": [{"operation": "add", "amount": 3}]},
        ]},
        status="pending_review",
    )
    pending = _number_line_scene(tmp_path).model_copy(
        update={"scene_id": "s2", "candidate_id": "c2"}
    )
    _seed_scene(client, already_chained)
    _seed_scene(client, pending)
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})
    assert resp.status_code == 400


def test_chain_rejects_mismatched_templates(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    number_line = _number_line_scene(tmp_path)
    array_grid = Scene(
        scene_id="s2", candidate_id="c2", template=static_ref(TemplateName.ARRAY_GRID),
        grade_level=1, params={"rows": 2, "cols": 3}, status="pending_review",
    )
    _seed_scene(client, number_line)
    _seed_scene(client, array_grid)
    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})
    assert resp.status_code == 400


def test_chain_rejects_text_card_template(tmp_path):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    scene1 = Scene(
        scene_id="s1", candidate_id="c1", template=static_ref(TemplateName.TEXT_CARD),
        grade_level=1, params={"headline": "x", "lines": ["y"]},
        status="pending_review",
    )
    scene2 = Scene(
        scene_id="s2", candidate_id="c2", template=static_ref(TemplateName.TEXT_CARD),
        grade_level=1, params={"headline": "x", "lines": ["y"]},
        status="pending_review",
    )
    _seed_scene(client, scene1)
    _seed_scene(client, scene2)

    resp = client.post("/storyboard/chain", json={"scene_ids": ["s1", "s2"]})
    assert resp.status_code == 400
    assert "text_card" in resp.json()["detail"].lower()


def test_chain_rejects_reusing_absorbed_constituent_scene_id(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), _candidate("c3")])
    a = _number_line_scene(tmp_path).model_copy(update={"scene_id": "a", "candidate_id": "c1"})
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    c = _number_line_scene(tmp_path).model_copy(update={"scene_id": "c", "candidate_id": "c3"})
    _seed_scene(client, a)
    _seed_scene(client, b)
    _seed_scene(client, c)

    with patch("app.routes.render_chained_scene_thumbnail"):
        first = client.post("/storyboard/chain", json={"scene_ids": ["a", "b"]})
    assert first.status_code == 200

    # "a" was absorbed into the chained scene above and removed from scene_order,
    # but its Scene object is deliberately kept around for ungroup. Re-submitting
    # it here must be rejected cleanly, not crash with an unhandled ValueError.
    resp = client.post("/storyboard/chain", json={"scene_ids": ["a", "c"]})

    assert resp.status_code == 400


def test_chain_splices_into_earliest_screen_position(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), _candidate("c3")])
    a = _number_line_scene(tmp_path).model_copy(update={"scene_id": "a", "candidate_id": "c1"})
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    c = _number_line_scene(tmp_path).model_copy(update={"scene_id": "c", "candidate_id": "c3"})
    _seed_scene(client, a)
    _seed_scene(client, b)
    _seed_scene(client, c)

    with patch("app.routes.render_chained_scene_thumbnail"):
        resp = client.post("/storyboard/chain", json={"scene_ids": ["a", "c"]})

    assert resp.status_code == 200
    new_id = resp.json()["scene_id"]

    from app.routes import store
    session = store.get(client.cookies.get("session_id"))
    assert session.scene_order == [new_id, "b"]


def test_chain_preserves_request_order_for_content_and_screen_order_for_restoration(tmp_path):
    client = _client()
    c3_candidate = _candidate("c3").model_copy(
        update={
            "source_excerpt": "Nine minus two.",
            "one_line_summary": "Detected: 9 - 2",
        }
    )
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), c3_candidate])
    a = _number_line_scene(tmp_path).model_copy(
        update={
            "scene_id": "a",
            "candidate_id": "c1",
            "params": {"start": 1, "steps": [{"operation": "add", "amount": 1}]},
        }
    )
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    c = _number_line_scene(tmp_path).model_copy(
        update={
            "scene_id": "c",
            "candidate_id": "c3",
            "grade_level": 6,
            "grade_overridden": True,
            "params": {"start": 9, "steps": [{"operation": "subtract", "amount": 2}]},
        }
    )
    _seed_scene(client, a)
    _seed_scene(client, b)
    _seed_scene(client, c)

    with patch("app.routes.render_chained_scene_thumbnail"):
        resp = client.post("/storyboard/chain", json={"scene_ids": ["c", "a"]})

    assert resp.status_code == 200
    body = resp.json()
    assert body["candidate_ids"] == ["c3", "c1"]
    assert [item["start"] for item in body["params"]["items"]] == [9, 1]
    assert body["source_excerpt"] == (
        "Nine minus two. / Sarah has 4 apples and buys 3 more."
    )
    assert body["detected_summary"] == "Detected: 9 - 2 / Detected: 4 + 3"
    assert body["grade_level"] == 6
    assert body["grade_overridden"] is True

    from app.routes import store
    session = store.get(client.cookies.get("session_id"))
    assert session.scene_order == [body["scene_id"], "b"]
    assert session.scene_chain_members[body["scene_id"]] == ["a", "c"]


def test_ungroup_restores_original_scenes_at_same_position(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), _candidate("c3")])
    a = _number_line_scene(tmp_path).model_copy(update={"scene_id": "a", "candidate_id": "c1"})
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    c = _number_line_scene(tmp_path).model_copy(update={"scene_id": "c", "candidate_id": "c3"})
    _seed_scene(client, a)
    _seed_scene(client, b)
    _seed_scene(client, c)

    with patch("app.routes.render_chained_scene_thumbnail"):
        chain_resp = client.post("/storyboard/chain", json={"scene_ids": ["a", "c"]})
    new_id = chain_resp.json()["scene_id"]

    resp = client.post(f"/storyboard/{new_id}/ungroup")

    assert resp.status_code == 200
    scenes = resp.json()["scenes"]
    assert [s["scene_id"] for s in scenes] == ["a", "c"]
    assert scenes[0]["candidate_id"] == "c1"
    assert scenes[1]["candidate_id"] == "c3"

    from app.routes import store
    session = store.get(client.cookies.get("session_id"))
    assert session.scene_order == ["a", "c", "b"]
    assert new_id not in session.scenes
    assert new_id not in session.scene_chain_members


def test_ungroup_restores_true_screen_order_when_request_order_is_reversed(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2"), _candidate("c3")])
    a = _number_line_scene(tmp_path).model_copy(update={"scene_id": "a", "candidate_id": "c1"})
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    c = _number_line_scene(tmp_path).model_copy(update={"scene_id": "c", "candidate_id": "c3"})
    _seed_scene(client, a)
    _seed_scene(client, b)
    _seed_scene(client, c)

    # Request lists scene_ids in reversed on-screen order (c, a instead of a, c).
    with patch("app.routes.render_chained_scene_thumbnail"):
        chain_resp = client.post("/storyboard/chain", json={"scene_ids": ["c", "a"]})
    new_id = chain_resp.json()["scene_id"]

    resp = client.post(f"/storyboard/{new_id}/ungroup")

    assert resp.status_code == 200
    scenes = resp.json()["scenes"]
    # Restored in true on-screen order (a before c), not request order (c before a).
    assert [s["scene_id"] for s in scenes] == ["a", "c"]
    assert scenes[0]["candidate_id"] == "c1"
    assert scenes[1]["candidate_id"] == "c3"

    from app.routes import store
    session = store.get(client.cookies.get("session_id"))
    assert session.scene_order == ["a", "c", "b"]
    assert new_id not in session.scenes
    assert new_id not in session.scene_chain_members


@pytest.mark.parametrize(
    ("method", "path", "body"),
    [
        ("patch", "/storyboard/a", {"grade_level": 4}),
        ("post", "/storyboard/a/approve", None),
        ("post", "/storyboard/a/reject", None),
        ("post", "/storyboard/a/retry", None),
    ],
)
def test_absorbed_scene_cannot_be_mutated(method, path, body, tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    a = _number_line_scene(tmp_path).model_copy(update={"scene_id": "a", "candidate_id": "c1"})
    b = _number_line_scene(tmp_path).model_copy(update={"scene_id": "b", "candidate_id": "c2"})
    _seed_scene(client, a, template="number_line")
    _seed_scene(client, b, template="number_line")

    with patch("app.routes.render_chained_scene_thumbnail"):
        chain_resp = client.post("/storyboard/chain", json={"scene_ids": ["a", "b"]})
    assert chain_resp.status_code == 200

    resp = client.request(method, path, json=body)

    assert resp.status_code == 404


def test_ungroup_unknown_or_non_chain_scene_is_404(tmp_path):
    client = _client()
    _upload_candidate(client)
    _seed_scene(client, _number_line_scene(tmp_path))

    resp = client.post("/storyboard/s1/ungroup")
    assert resp.status_code == 404

    resp = client.post("/storyboard/nope/ungroup")
    assert resp.status_code == 404


def _chained_number_line_scene(candidate_ids=("c1", "c2")):
    from app.models.scene import Scene, TemplateName
    from app.templates.registry import static_ref

    return Scene(
        scene_id="s1",
        candidate_ids=list(candidate_ids),
        template=static_ref(TemplateName.NUMBER_LINE),
        grade_level=1,
        params={"items": [
            {"start": 4, "steps": [{"operation": "add", "amount": 3}]},
            {"start": 5, "steps": [{"operation": "subtract", "amount": 1}]},
        ]},
        status="pending_review",
    )


def test_patch_chained_scene_valid_item_edit_re_renders_thumbnail():
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    _seed_scene(client, _chained_number_line_scene())

    with patch("app.routes.render_chained_scene_thumbnail") as thumb:
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"items": [
                {"start": 10, "steps": [{"operation": "add", "amount": 1}]},
                {"start": 5, "steps": [{"operation": "subtract", "amount": 1}]},
            ]}},
        )

    assert resp.status_code == 200
    assert resp.json()["params"]["items"][0]["start"] == 10
    thumb.assert_called_once()


def test_patch_chained_scene_invalid_item_returns_422():
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    _seed_scene(client, _chained_number_line_scene())

    with patch("app.routes.render_chained_scene_thumbnail") as thumb:
        resp = client.patch(
            "/storyboard/s1",
            json={"params": {"items": [
                # start=1 then subtract 5 -> running total goes negative -> guard rejects.
                {"start": 1, "steps": [{"operation": "subtract", "amount": 5}]},
                {"start": 5, "steps": [{"operation": "subtract", "amount": 1}]},
            ]}},
        )

    assert resp.status_code == 422
    assert resp.json()["detail"]["errors"]
    thumb.assert_not_called()


def test_render_chained_scene_uses_chained_render_path(tmp_path):
    client = _client()
    _upload_candidates(client, [_candidate("c1"), _candidate("c2")])
    chained = _chained_number_line_scene().model_copy(update={"status": "approved"})
    _seed_scene(client, chained)

    def fake_render(template, params, out):
        out.write_bytes(b"mp4")
        return out

    with patch(
        "app.routes.render_chained_scene_to_mp4", side_effect=fake_render
    ) as chained_render, patch("app.routes.render_scene_to_mp4") as solo_render:
        resp = client.post("/render")

    assert resp.status_code == 200
    clips = resp.json()["clips"]
    assert clips[0]["candidate_ids"] == ["c1", "c2"]
    assert clips[0]["clip_url"].startswith("/clips/")
    chained_render.assert_called_once()
    solo_render.assert_not_called()


def test_options_loads_the_snapshot_for_the_requesting_session():
    """A teacher's own approved templates must reach their own /options.

    load_enabled_snapshot filters by owner, so passing the session id is what
    makes a privately-approved template visible to the session that approved it
    and invisible to everyone else.
    """
    from app.config import get_settings
    from app.meta.dynamic_templates import EnabledSnapshot

    client = _client()
    _upload_candidate(client)
    session_id = client.cookies.get("session_id")

    settings = get_settings()
    settings.meta_dynamic_classifier_enabled = True
    try:
        with patch("app.routes.meta_session") as mock_meta_session, patch(
            "app.routes.load_enabled_snapshot", return_value=EnabledSnapshot(_entries={})
        ) as mock_load_snapshot, patch(
            "app.routes.classify_candidate", return_value=_classification()
        ):
            mock_meta_session.return_value.__enter__.return_value = object()
            client.post("/options", json={"candidate_ids": ["c1"]})
    finally:
        settings.meta_dynamic_classifier_enabled = False

    assert mock_load_snapshot.call_args.kwargs["owner_session_id"] == session_id
