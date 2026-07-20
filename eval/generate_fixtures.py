from pathlib import Path

from pptx import Presentation

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def _new_presentation():
    presentation = Presentation()
    return presentation, presentation.slide_layouts[1]


def build_zero_candidate_deck(path: Path) -> None:
    presentation, layout = _new_presentation()

    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "Welcome"
    slide1.placeholders[1].text = "Chapter 4: Fractions Overview"

    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Agenda"
    slide2.placeholders[1].text = "Standards: 3.OA.A.1\nHomework due Friday"

    presentation.save(path)


def build_distractor_heavy_deck(path: Path) -> None:
    presentation, layout = _new_presentation()

    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Warm Up"
    slide.placeholders[1].text = (
        "Date: 3/14. Standard 3.OA.A.1. Page 42.\n"
        "Sarah has 4 apples and buys 3 more apples. How many apples does she have now?"
    )

    presentation.save(path)


def build_ambiguous_phrasing_deck(path: Path) -> None:
    presentation, layout = _new_presentation()

    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Think About It"
    slide.placeholders[1].text = (
        "There are some red apples and some green apples. There are 4 more red "
        "apples than green apples. How many apples are there in all?"
    )

    presentation.save(path)


def main() -> None:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    build_zero_candidate_deck(FIXTURES_DIR / "zero_candidate_deck.pptx")
    build_distractor_heavy_deck(FIXTURES_DIR / "distractor_heavy_deck.pptx")
    build_ambiguous_phrasing_deck(FIXTURES_DIR / "ambiguous_phrasing_deck.pptx")
    print(f"Wrote fixtures to {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
